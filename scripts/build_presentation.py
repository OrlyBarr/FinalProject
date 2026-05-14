#!/usr/bin/env python3
"""
Israel Public Transit Monitoring — Final Project Presentation Builder
Creates a comprehensive PPTX presentation with project architecture.
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from copy import deepcopy
from lxml import etree

# ─────────────────────────────────────────────────────────────────────────────
# Color Palette — light teal / sky blue
# ─────────────────────────────────────────────────────────────────────────────
TEAL_DARK    = RGBColor(0x0B, 0x6E, 0x86)   # primary text
TEAL_MID     = RGBColor(0x14, 0xA5, 0xB8)   # secondary
TEAL_LIGHT   = RGBColor(0x7D, 0xD3, 0xDC)   # backgrounds
SKY_BLUE     = RGBColor(0xA6, 0xE3, 0xE9)   # accent / boxes
ICE_BG       = RGBColor(0xEA, 0xF8, 0xFB)   # slide background
WHITE        = RGBColor(0xFF, 0xFF, 0xFF)
DARK_TXT     = RGBColor(0x1B, 0x2C, 0x3A)   # body text
GRAY         = RGBColor(0x6B, 0x7C, 0x8A)
ACCENT_GREEN = RGBColor(0x4F, 0xB7, 0x82)
ACCENT_AMBER = RGBColor(0xE8, 0xA8, 0x4F)
RED          = RGBColor(0xD4, 0x5D, 0x5D)


def hebrew_run(p, text, *, size=18, bold=False, color=DARK_TXT, font='Arial'):
    """Add a Hebrew text run with RTL support."""
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return run


def set_rtl(paragraph):
    """Set paragraph direction to RTL for Hebrew."""
    pPr = paragraph._pPr or paragraph._p.get_or_add_pPr()
    pPr.set('rtl', '1')


def add_textbox(slide, left, top, width, height, *, fill=None, line=None,
                shadow=False):
    """Add a positioned textbox / shape and return it."""
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                 left, top, width, height)
    if fill is None:
        box.fill.background()
    else:
        box.fill.solid()
        box.fill.fore_color.rgb = fill
    if line is None:
        box.line.fill.background()
    else:
        box.line.color.rgb = line
        box.line.width = Pt(1)
    if not shadow:
        # Disable default shadow
        spPr = box._element.spPr
        for el in spPr.findall(qn('a:effectLst')):
            spPr.remove(el)
        effect = etree.SubElement(spPr, qn('a:effectLst'))
    box.text_frame.margin_left   = Inches(0.15)
    box.text_frame.margin_right  = Inches(0.15)
    box.text_frame.margin_top    = Inches(0.08)
    box.text_frame.margin_bottom = Inches(0.08)
    box.text_frame.word_wrap = True
    return box


def add_text(slide, left, top, width, height, text, *,
             size=18, bold=False, color=DARK_TXT, align=PP_ALIGN.RIGHT,
             rtl=True, font='Arial'):
    """Plain textbox helper."""
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(0.05)
    tf.margin_top  = tf.margin_bottom = Inches(0.02)
    p = tf.paragraphs[0]
    p.alignment = align
    if rtl:
        set_rtl(p)
    hebrew_run(p, text, size=size, bold=bold, color=color, font=font)
    return tb


def slide_bg(slide, color=ICE_BG):
    """Set slide background color."""
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0,
                                Emu(12192000), Emu(6858000))  # 16:9 size
    bg.fill.solid()
    bg.fill.fore_color.rgb = color
    bg.line.fill.background()
    bg.shadow.inherit = False
    # Send to back
    spTree = bg._element.getparent()
    spTree.remove(bg._element)
    spTree.insert(2, bg._element)


def add_header(slide, title_he, subtitle_he=None, slide_num=None, total=None):
    """Add page header — title bar + slide number."""
    # Top accent bar
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                 0, 0, Emu(12192000), Inches(0.12))
    bar.fill.solid()
    bar.fill.fore_color.rgb = TEAL_MID
    bar.line.fill.background()

    # Title
    add_text(slide, Inches(0.4), Inches(0.25), Inches(12.4), Inches(0.6),
             title_he, size=30, bold=True, color=TEAL_DARK)

    if subtitle_he:
        add_text(slide, Inches(0.4), Inches(0.85), Inches(12.4), Inches(0.4),
                 subtitle_he, size=15, color=GRAY)

    if slide_num is not None and total:
        add_text(slide, Inches(11.6), Inches(7.0), Inches(1.2), Inches(0.3),
                 f"{slide_num} / {total}", size=11, color=GRAY,
                 align=PP_ALIGN.RIGHT, rtl=False)


# ═══════════════════════════════════════════════════════════════════════════
# Build the presentation
# ═══════════════════════════════════════════════════════════════════════════
prs = Presentation()
prs.slide_width  = Inches(13.333)   # 16:9 widescreen
prs.slide_height = Inches(7.5)

blank_layout = prs.slide_layouts[6]
TOTAL = 15   # we'll create this many


# ─────────── SLIDE 1: Cover ───────────
s = prs.slides.add_slide(blank_layout)
slide_bg(s, ICE_BG)

# Big colored band
band = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
    0, Inches(2.0), Emu(12192000), Inches(3.5))
band.fill.solid()
band.fill.fore_color.rgb = TEAL_LIGHT
band.line.fill.background()

# Decorative circles
for x, y, r, c in [(0.8, 0.6, 0.6, SKY_BLUE), (12.0, 0.4, 0.4, TEAL_MID),
                   (11.5, 6.5, 0.7, SKY_BLUE), (0.4, 6.6, 0.5, TEAL_MID)]:
    c1 = s.shapes.add_shape(MSO_SHAPE.OVAL,
        Inches(x - r), Inches(y - r), Inches(2*r), Inches(2*r))
    c1.fill.solid(); c1.fill.fore_color.rgb = c
    c1.line.fill.background()

# Title
add_text(s, Inches(1), Inches(2.3), Inches(11.3), Inches(1.3),
         "🚌 ניטור תחבורה ציבורית — זמן אמת",
         size=54, bold=True, color=TEAL_DARK, align=PP_ALIGN.CENTER)

# Subtitle Hebrew
add_text(s, Inches(1), Inches(3.5), Inches(11.3), Inches(0.7),
         "פלטפורמת ניטור Big Data לתחבורה ציבורית בגוש דן ותל אביב",
         size=22, color=DARK_TXT, align=PP_ALIGN.CENTER)

# Subtitle English
add_text(s, Inches(1), Inches(4.2), Inches(11.3), Inches(0.5),
         "Real-Time Public Transit Monitoring — Big Data Pipeline",
         size=16, color=GRAY, align=PP_ALIGN.CENTER, rtl=False, font='Arial')

# Footer info
add_text(s, Inches(1), Inches(6.3), Inches(11.3), Inches(0.5),
         "פרויקט סיום • מאי 2026",
         size=18, bold=True, color=TEAL_DARK, align=PP_ALIGN.CENTER)

add_text(s, Inches(1), Inches(6.8), Inches(11.3), Inches(0.4),
         "Kafka • Airflow • Elasticsearch • MinIO • Kibana • Flask",
         size=13, color=GRAY, align=PP_ALIGN.CENTER, rtl=False, font='Consolas')


# ─────────── SLIDE 2: Project Overview ───────────
s = prs.slides.add_slide(blank_layout); slide_bg(s)
add_header(s, "סקירה כללית של הפרויקט",
           "מה המערכת עושה?", 2, TOTAL)

# Left big text
add_text(s, Inches(0.5), Inches(1.5), Inches(12.3), Inches(0.6),
         "פלטפורמה לניטור תחבורה ציבורית בזמן אמת",
         size=28, bold=True, color=TEAL_DARK)

# Description
desc = ("המערכת אוספת נתונים בזמן אמת מכל קווי האוטובוס והרכבת בגוש דן, "
        "מעבדת אותם דרך תשתית Big Data, ומציגה אותם בממשק אינטראקטיבי בסגנון Moovit "
        "וכן בדשבורדים אנליטיים ב-Kibana.")
add_text(s, Inches(0.5), Inches(2.2), Inches(12.3), Inches(0.9),
         desc, size=16, color=DARK_TXT)

# 4 feature boxes
features = [
    ("🚌", "מעקב חי", "מיקום GPS של אלפי אוטובוסים ורכבות מתעדכן כל 5 דקות"),
    ("📊", "ניתוח נתונים", "דשבורדים ב-Kibana — איחורים, עומסי תנועה, פעילות מפעילים"),
    ("🚦", "תנועה בכבישים", "מהירויות ועומסים מ-HERE Traffic API לכל הארץ"),
    ("🤖", "ממשק חכם", "אפליקציית web בסגנון Moovit + AI agent עם Gemini"),
]
for i, (icon, title, desc) in enumerate(features):
    col = i % 2
    row = i // 2
    x = Inches(0.5 + col * 6.4)
    y = Inches(3.4 + row * 1.95)
    box = add_textbox(s, x, y, Inches(6.0), Inches(1.7),
                      fill=WHITE, line=TEAL_LIGHT)
    # Icon
    add_text(s, x + Inches(0.2), y + Inches(0.15), Inches(1.0), Inches(0.7),
             icon, size=36, color=TEAL_MID, align=PP_ALIGN.CENTER, rtl=False)
    # Title
    add_text(s, x + Inches(1.2), y + Inches(0.2), Inches(4.6), Inches(0.5),
             title, size=20, bold=True, color=TEAL_DARK)
    # Desc
    add_text(s, x + Inches(1.2), y + Inches(0.75), Inches(4.6), Inches(0.85),
             desc, size=13, color=DARK_TXT)


# ─────────── SLIDE 3: High-Level Architecture ───────────
s = prs.slides.add_slide(blank_layout); slide_bg(s)
add_header(s, "ארכיטקטורה — תמונה כוללת",
           "ארבע שכבות: מקור נתונים → ייבוא → אחסון → הצגה", 3, TOTAL)

# 4 horizontal layers
layers = [
    ("מקור הנתונים  📡", "External APIs",
     ["Open Bus Stride (Hasadna)", "MOT GTFS-RT", "HERE Traffic API"],
     SKY_BLUE),
    ("שכבת ייבוא  🔄", "Ingestion Layer",
     ["Apache Kafka — 9 topics", "Apache Airflow — orchestration",
      "Python Producers"],
     TEAL_LIGHT),
    ("שכבת אחסון  💾", "Storage Layer",
     ["MinIO — Data Lake (raw JSON)", "Elasticsearch — searchable index",
      "PostgreSQL — GTFS static"],
     RGBColor(0xC8, 0xEC, 0xF0)),
    ("שכבת הצגה  🖥️", "Presentation Layer",
     ["Flask Bot (port 5000)", "Kibana Dashboards", "Moovit-style web app"],
     RGBColor(0xB0, 0xE0, 0xE6)),
]

y = Inches(1.45)
h = Inches(1.32)
for label_he, label_en, items, color in layers:
    # Layer box
    box = add_textbox(s, Inches(0.4), y, Inches(12.5), h,
                      fill=color, line=TEAL_MID)
    # Hebrew label (right side)
    add_text(s, Inches(9.5), y + Inches(0.1), Inches(3.3), Inches(0.5),
             label_he, size=20, bold=True, color=TEAL_DARK)
    add_text(s, Inches(9.5), y + Inches(0.65), Inches(3.3), Inches(0.4),
             label_en, size=12, color=GRAY, font='Arial', rtl=False)
    # Items (3 columns)
    for i, item in enumerate(items):
        item_box = add_textbox(s, Inches(0.6 + i * 2.95), y + Inches(0.32),
                               Inches(2.85), Inches(0.7),
                               fill=WHITE)
        add_text(s, Inches(0.6 + i * 2.95), y + Inches(0.45),
                 Inches(2.85), Inches(0.55),
                 item, size=12, color=DARK_TXT,
                 align=PP_ALIGN.CENTER, rtl=False)
    # Down arrow between layers (except after last)
    if y < Inches(5.0):
        arrow = s.shapes.add_shape(MSO_SHAPE.DOWN_ARROW,
            Inches(6.4), y + h + Inches(0.02), Inches(0.5), Inches(0.13))
        arrow.fill.solid(); arrow.fill.fore_color.rgb = TEAL_DARK
        arrow.line.fill.background()
    y += h + Inches(0.18)


# ─────────── SLIDE 4: Data Sources ───────────
s = prs.slides.add_slide(blank_layout); slide_bg(s)
add_header(s, "מקורות הנתונים  📡",
           "3 APIs ציבוריים — ללא צורך ברישום", 4, TOTAL)

sources = [
    ("Open Bus Stride", "open-bus-stride-api.hasadna.org.il",
     "API פתוח של עמותת \"הסדנא לידע ציבורי\". מקור הנתונים העיקרי.",
     ["מיקומי GPS של אוטובוסים ורכבות (SIRI)",
      "מסלולי קווים מלאים (GTFS)",
      "לוחות זמנים מתוכננים",
      "12 מפעילים: דן, אגד, סופרבוס, נתיב אקספרס, רכבת ישראל..."],
     TEAL_MID),
    ("MOT GTFS-RT", "gtfs.mot.gov.il",
     "פיד פרוטוקול-באפר רשמי ממשרד התחבורה.",
     ["TripUpdates — עיכובי נסיעות",
      "ServiceAlerts — התראות שירות",
      "VehiclePositions — מיקומי רכבים",
      "⚠️ חסום ע\"י WAF — מתוקן דרך Stride"],
     TEAL_DARK),
    ("HERE Traffic v7", "data.traffic.hereapi.com",
     "נתוני זרימת תנועה לכל מקטע כביש.",
     ["מהירות בפועל ומהירות חופשית",
      "Jam Factor (0-10) — מקדם עומס",
      "11 אזורי כיסוי בישראל",
      "טווח שינוי: כל 5 דקות"],
     ACCENT_GREEN),
]

y = Inches(1.5)
for name, url, desc, bullets, color in sources:
    box = add_textbox(s, Inches(0.4), y, Inches(12.5), Inches(1.75),
                      fill=WHITE, line=color)
    # Color stripe on the right (Hebrew side)
    stripe = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
        Inches(12.6), y, Inches(0.3), Inches(1.75))
    stripe.fill.solid(); stripe.fill.fore_color.rgb = color
    stripe.line.fill.background()
    # Title + URL
    add_text(s, Inches(0.6), y + Inches(0.1), Inches(11.7), Inches(0.4),
             name, size=18, bold=True, color=color, align=PP_ALIGN.LEFT,
             rtl=False, font='Arial')
    add_text(s, Inches(0.6), y + Inches(0.45), Inches(11.7), Inches(0.3),
             url, size=11, color=GRAY, align=PP_ALIGN.LEFT,
             rtl=False, font='Consolas')
    # Description (Hebrew)
    add_text(s, Inches(0.6), y + Inches(0.75), Inches(11.9), Inches(0.3),
             desc, size=12, color=DARK_TXT)
    # Bullets - 4 in 2 rows of 2
    for i, b in enumerate(bullets):
        col = i % 2
        row = i // 2
        bx = Inches(0.6 + col * 6.05)
        by = y + Inches(1.10 + row * 0.30)
        add_text(s, bx, by, Inches(6.0), Inches(0.28),
                 "• " + b, size=11, color=DARK_TXT)
    y += Inches(1.85)


# ─────────── SLIDE 5: Ingestion (Kafka + Airflow) ───────────
s = prs.slides.add_slide(blank_layout); slide_bg(s)
add_header(s, "שכבת הייבוא  🔄",
           "Apache Kafka + Apache Airflow", 5, TOTAL)

# Left: Kafka
kbox = add_textbox(s, Inches(0.4), Inches(1.4), Inches(6.2), Inches(5.5),
                   fill=WHITE, line=TEAL_MID)
add_text(s, Inches(0.6), Inches(1.55), Inches(5.8), Inches(0.5),
         "Apache Kafka", size=22, bold=True, color=TEAL_DARK,
         align=PP_ALIGN.LEFT, rtl=False)
add_text(s, Inches(0.6), Inches(2.0), Inches(5.8), Inches(0.4),
         "תור הודעות מבוזר — מאחסן את כל הזרמים בזמן אמת", size=12,
         color=DARK_TXT)

topics = [
    ("bus-positions", "4 partitions", "500 הודעות/דקה"),
    ("train-positions", "4 partitions", "300 הודעות/דקה"),
    ("traffic-data", "2 partitions", "443 הודעות/דקה"),
    ("trip-updates", "2 partitions", "500 הודעות/דקה"),
    ("service-alerts", "2 partitions", "133 הודעות/דקה"),
    ("delay-events", "1 partition", "אירועי איחור"),
    ("pipeline-errors", "1 partition", "שגיאות"),
]
ty = Inches(2.55)
for t, p_, throughput in topics:
    add_text(s, Inches(0.7), ty, Inches(2.5), Inches(0.32),
             "📨 " + t, size=11, color=TEAL_DARK,
             align=PP_ALIGN.LEFT, rtl=False, font='Consolas')
    add_text(s, Inches(3.2), ty, Inches(1.4), Inches(0.32),
             p_, size=10, color=GRAY,
             align=PP_ALIGN.LEFT, rtl=False)
    add_text(s, Inches(4.6), ty, Inches(1.9), Inches(0.32),
             throughput, size=10, color=DARK_TXT)
    ty += Inches(0.36)

# Right: Airflow
abox = add_textbox(s, Inches(6.8), Inches(1.4), Inches(6.2), Inches(5.5),
                   fill=WHITE, line=TEAL_MID)
add_text(s, Inches(7.0), Inches(1.55), Inches(5.8), Inches(0.5),
         "Apache Airflow", size=22, bold=True, color=TEAL_DARK,
         align=PP_ALIGN.LEFT, rtl=False)
add_text(s, Inches(7.0), Inches(2.0), Inches(5.8), Inches(0.4),
         "מתזמן (orchestrator) — מריץ DAGs בלוחות זמנים", size=12,
         color=DARK_TXT)

dags = [
    ("dag_realtime_ingestion", "כל 2 דקות", "מריץ Producers"),
    ("dag_etl_transform", "כל 5 דקות", "מעבד Kafka → Delay events"),
    ("dag_es_indexer", "כל 5 דקות", "MinIO → Elasticsearch"),
    ("dag_direct_traffic", "כל 5 דקות", "HERE → MinIO ישיר"),
    ("dag_direct_alerts", "כל שעה", "התראות שירות"),
    ("dag_gtfs_load", "פעם ביום", "טעינת GTFS סטטי"),
]
dy = Inches(2.55)
for name, schedule, desc in dags:
    add_text(s, Inches(7.1), dy, Inches(2.9), Inches(0.32),
             "⚙️ " + name, size=11, color=TEAL_DARK,
             align=PP_ALIGN.LEFT, rtl=False, font='Consolas')
    add_text(s, Inches(10.0), dy, Inches(1.3), Inches(0.32),
             schedule, size=10, color=ACCENT_AMBER, bold=True,
             align=PP_ALIGN.LEFT)
    add_text(s, Inches(7.1), dy + Inches(0.30), Inches(5.7), Inches(0.28),
             "    " + desc, size=10, color=GRAY,
             align=PP_ALIGN.LEFT)
    dy += Inches(0.66)


# ─────────── SLIDE 6: Storage Layer ───────────
s = prs.slides.add_slide(blank_layout); slide_bg(s)
add_header(s, "שכבת האחסון  💾",
           "שלוש מערכות אחסון לפי סוג השימוש", 6, TOTAL)

storages = [
    ("MinIO", "S3-compatible Object Store",
     "אגם נתונים — JSON גולמי",
     [("Bucket:", "israel-transit-lake"),
      ("נתיב:", "raw/{topic}/year=Y/month=M/day=D/hour=H"),
      ("גרסאות:", "raw + processed"),
      ("נפח כיום:", "~500 MB יומי")],
     TEAL_MID, "🗄️"),
    ("Elasticsearch 8.8", "Distributed Search Engine",
     "אינדקס לחיפוש ושאילתות מהירות",
     [("transit-bus-positions:", "268,771 docs"),
      ("transit-train-positions:", "508,772 docs"),
      ("transit-traffic:", "67,423 docs"),
      ("transit-trip-updates:", "51,874 docs"),
      ("transit-service-alerts:", "17,512 docs"),
      ],
     TEAL_DARK, "🔍"),
    ("PostgreSQL 13", "Relational Database",
     "GTFS סטטי + מטא-דאטה של Airflow",
     [("Schema:", "gtfs (routes, stops, trips, stop_times)"),
      ("Schema:", "airflow (DAG runs, tasks)"),
      ("Port:", "15432 (כדי לא להתנגש עם המערכת)")],
     ACCENT_GREEN, "🐘"),
]

y = Inches(1.4)
for name, en, role, items, color, icon in storages:
    h = Inches(1.6 if len(items) <= 4 else 1.85)
    box = add_textbox(s, Inches(0.4), y, Inches(12.5), h,
                      fill=WHITE, line=color)
    # Icon big on right
    add_text(s, Inches(11.8), y + Inches(0.2), Inches(1.0), Inches(1.2),
             icon, size=48, color=color, align=PP_ALIGN.CENTER, rtl=False)
    # Title
    add_text(s, Inches(0.6), y + Inches(0.15), Inches(8), Inches(0.4),
             name, size=20, bold=True, color=color,
             align=PP_ALIGN.LEFT, rtl=False)
    add_text(s, Inches(0.6), y + Inches(0.50), Inches(8), Inches(0.3),
             en, size=11, color=GRAY,
             align=PP_ALIGN.LEFT, rtl=False, font='Arial')
    # Role (Hebrew)
    add_text(s, Inches(2.5), y + Inches(0.15), Inches(8.8), Inches(0.4),
             role, size=14, color=DARK_TXT)
    # Items in 2 columns
    for i, (k, v) in enumerate(items):
        col = i % 2
        row = i // 2
        bx = Inches(0.6 + col * 5.6)
        by = y + Inches(0.85 + row * 0.30)
        add_text(s, bx, by, Inches(2.0), Inches(0.28),
                 k, size=11, bold=True, color=DARK_TXT,
                 align=PP_ALIGN.LEFT, rtl=False)
        add_text(s, bx + Inches(2.0), by, Inches(3.5), Inches(0.28),
                 v, size=11, color=GRAY,
                 align=PP_ALIGN.LEFT, rtl=False, font='Consolas')
    y += h + Inches(0.18)


# ─────────── SLIDE 7: Data Flow Diagram ───────────
s = prs.slides.add_slide(blank_layout); slide_bg(s)
add_header(s, "זרימת נתונים מקצה לקצה",
           "מסע של רשומה אחת מה-API ועד למשתמש", 7, TOTAL)

# 6 boxes left-to-right with arrows
steps = [
    ("1️⃣", "API חיצוני",  "Stride / HERE",   TEAL_MID),
    ("2️⃣", "Producer",    "Python (asyncio)", TEAL_DARK),
    ("3️⃣", "Kafka",       "Topic + Partitions", TEAL_MID),
    ("4️⃣", "Consumer",    "Airflow Task",     TEAL_DARK),
    ("5️⃣", "MinIO + ES",  "Storage",          ACCENT_GREEN),
    ("6️⃣", "Frontend",    "Flask / Kibana",   TEAL_MID),
]

box_w = Inches(1.85)
gap   = Inches(0.18)
total_w = 6 * box_w + 5 * gap
start_x = (Inches(13.333) - total_w) / 2
y_box = Inches(2.3)
h_box = Inches(2.0)

for i, (num, title, sub, color) in enumerate(steps):
    x = start_x + i * (box_w + gap)
    # Box
    box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                             x, y_box, box_w, h_box)
    box.fill.solid(); box.fill.fore_color.rgb = color
    box.line.fill.background()
    # Number
    add_text(s, x, y_box + Inches(0.15), box_w, Inches(0.5),
             num, size=24, color=WHITE, align=PP_ALIGN.CENTER, rtl=False)
    # Title
    add_text(s, x, y_box + Inches(0.7), box_w, Inches(0.5),
             title, size=16, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER)
    # Subtitle
    add_text(s, x, y_box + Inches(1.25), box_w, Inches(0.55),
             sub, size=11, color=WHITE,
             align=PP_ALIGN.CENTER, rtl=False, font='Consolas')
    # Right-arrow between boxes (except after last)
    if i < 5:
        # In Hebrew RTL we still flow visually L→R for tech diagrams
        arrow = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,
            x + box_w + Inches(0.01), y_box + Inches(0.85),
            Inches(0.16), Inches(0.30))
        arrow.fill.solid(); arrow.fill.fore_color.rgb = TEAL_DARK
        arrow.line.fill.background()

# Bottom annotation
add_text(s, Inches(0.5), Inches(5.0), Inches(12.3), Inches(0.5),
         "💡 דוגמה: GPS של אוטובוס דן קו 1 ביציאה מתחנה",
         size=15, bold=True, color=TEAL_DARK, align=PP_ALIGN.CENTER)

example = ("Stride API → BusProducer → Kafka topic 'bus-positions' → "
           "Airflow consumer → MinIO (raw JSON + processed) → "
           "Elasticsearch (indexed) → המשתמש רואה במפה")
add_text(s, Inches(0.5), Inches(5.5), Inches(12.3), Inches(1.5),
         example, size=12, color=GRAY, align=PP_ALIGN.CENTER, rtl=False)


# ─────────── SLIDE 8: Resilience / Store-and-Forward ───────────
s = prs.slides.add_slide(blank_layout); slide_bg(s)
add_header(s, "עמידות המערכת  🛡️",
           "Store-and-Forward — איסוף ממשיך גם כשהקונטיינרים נופלים", 8, TOTAL)

# Top intro
add_text(s, Inches(0.5), Inches(1.4), Inches(12.3), Inches(0.5),
         "💡 הבעיה: מה קורה אם MinIO או Elasticsearch נופלים ל-30 דקות?",
         size=18, bold=True, color=TEAL_DARK)

add_text(s, Inches(0.5), Inches(1.95), Inches(12.3), Inches(0.55),
         "פתרון: ה-cron pipeline שלי רץ על ה-VM host ולא תלוי בקונטיינרים. "
         "אם MinIO לא זמין, הנתונים נשמרים ל-buffer מקומי ועולים אוטומטית כשהוא חוזר.",
         size=13, color=DARK_TXT)

# Three scenarios diagram
scenarios = [
    ("✅", "מצב רגיל",
     "API → Producer → MinIO → ES",
     "הנתונים זורמים ישירות, אפס פערים",
     ACCENT_GREEN),
    ("⚠️", "MinIO נופל",
     "API → Producer → 💾 Local Buffer",
     "נשמר על דיסק ב-/buffer/minio_pending/",
     ACCENT_AMBER),
    ("🔄", "MinIO חוזר",
     "Buffer → Flush → MinIO → ES",
     "אוטומטי בריצת cron הבאה (תוך 5 דק')",
     TEAL_MID),
]

w = Inches(4.15)
gap = Inches(0.15)
y = Inches(2.7)
h = Inches(2.6)
for i, (icon, title, flow, desc, color) in enumerate(scenarios):
    x = Inches(0.4 + i * (4.15 + 0.15))
    box = add_textbox(s, x, y, w, h, fill=WHITE, line=color)
    # Color band on top
    head = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, Inches(0.7))
    head.fill.solid(); head.fill.fore_color.rgb = color
    head.line.fill.background()
    # Icon + title in band
    add_text(s, x + Inches(0.15), y + Inches(0.1), Inches(0.7), Inches(0.5),
             icon, size=24, color=WHITE, align=PP_ALIGN.CENTER, rtl=False)
    add_text(s, x + Inches(0.85), y + Inches(0.15), w - Inches(1), Inches(0.45),
             title, size=18, bold=True, color=WHITE, align=PP_ALIGN.RIGHT)
    # Flow (monospace)
    add_text(s, x + Inches(0.15), y + Inches(0.85), w - Inches(0.3), Inches(0.5),
             flow, size=11, bold=True, color=color,
             align=PP_ALIGN.CENTER, rtl=False, font='Consolas')
    # Down arrow
    add_text(s, x + Inches(0.15), y + Inches(1.35), w - Inches(0.3), Inches(0.4),
             "↓", size=24, color=GRAY, align=PP_ALIGN.CENTER, rtl=False)
    # Description
    add_text(s, x + Inches(0.2), y + Inches(1.75), w - Inches(0.4), Inches(0.7),
             desc, size=12, color=DARK_TXT, align=PP_ALIGN.CENTER)

# Bottom: numbers
add_text(s, Inches(0.5), Inches(5.55), Inches(12.3), Inches(0.4),
         "🧪 נבדק מציאותית: זוהה outage → 0 איבוד נתונים → flush של 100% תוך הריצה הבאה",
         size=13, bold=True, color=ACCENT_GREEN, align=PP_ALIGN.CENTER)

# Tech detail
add_text(s, Inches(0.5), Inches(6.1), Inches(12.3), Inches(0.4),
         "Buffer מתמיד ב-/home/local_admin/finalproject/buffer/  (לא נמחק עם reboot)",
         size=11, color=GRAY, align=PP_ALIGN.CENTER, rtl=False, font='Consolas')

add_text(s, Inches(0.5), Inches(6.55), Inches(12.3), Inches(0.4),
         "📐 הדפוס נקרא Store-and-Forward — אותו דפוס ש-Kafka משתמש בו פנימית",
         size=12, color=TEAL_DARK, align=PP_ALIGN.CENTER)


# ─────────── SLIDE 9: Frontend ───────────
s = prs.slides.add_slide(blank_layout); slide_bg(s)
add_header(s, "ממשק המשתמש  🖥️",
           "Flask Bot + Moovit-style UI + Kibana", 9, TOTAL)

components = [
    ("🚌 Moovit-style App", "moovit.html",
     "ממשק משתמש בסגנון Moovit",
     ["מפה אינטראקטיבית (Leaflet)",
      "12 מפעילים — לחיצה מציגה קווים פעילים",
      "מסלול קו עם תחנות אמיתיות מ-GTFS",
      "סינון לקווים בתנועה (velocity > 0)"],
     TEAL_MID),
    ("🤖 Flask Bot API", "bot.py — port 5000",
     "שרת Python שמשרת את ה-UI ועושה proxy ל-APIs",
     ["/line_info — תחנות הקו (GTFS)",
      "/api/active_lines — קווים פעילים",
      "/proxy/stride/* — מעקף CORS",
      "/ask — Gemini AI agent"],
     TEAL_DARK),
    ("📊 Kibana Dashboards", "Kibana 8.8 — port 5601",
     "ויזואליזציה אנליטית של הנתונים",
     ["ניטור תחבורה ציבורית — זמן אמת",
      "ניתוח תנועה ועיכובים — גוש דן",
      "70+ ויזואליזציות שונות",
      "5 index patterns לחיפוש מתקדם"],
     ACCENT_GREEN),
]

w = Inches(4.15)
gap = Inches(0.15)
y = Inches(1.5)
h = Inches(5.4)
x = Inches(0.4)
for name, sub, desc, items, color in components:
    box = add_textbox(s, x, y, w, h, fill=WHITE, line=color)
    # Top color band
    band = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
        x, y, w, Inches(0.7))
    band.fill.solid(); band.fill.fore_color.rgb = color
    band.line.fill.background()
    add_text(s, x + Inches(0.15), y + Inches(0.1), w - Inches(0.3), Inches(0.55),
             name, size=18, bold=True, color=WHITE,
             align=PP_ALIGN.LEFT)
    # Subtitle
    add_text(s, x + Inches(0.15), y + Inches(0.85), w - Inches(0.3), Inches(0.3),
             sub, size=11, color=GRAY, font='Consolas',
             align=PP_ALIGN.LEFT, rtl=False)
    # Desc
    add_text(s, x + Inches(0.15), y + Inches(1.2), w - Inches(0.3), Inches(0.7),
             desc, size=12, color=DARK_TXT)
    # Items
    for i, item in enumerate(items):
        add_text(s, x + Inches(0.2), y + Inches(2.1 + i * 0.65),
                 w - Inches(0.4), Inches(0.6),
                 "• " + item, size=12, color=DARK_TXT)
    x += w + gap


# ─────────── SLIDE 10: Live Numbers ───────────
s = prs.slides.add_slide(blank_layout); slide_bg(s)
add_header(s, "המערכת בזמן אמת  📈",
           "מספרים מתוך ה-VM הפועל ברגע זה", 10, TOTAL)

# Big metric cards 3x3
metrics = [
    ("268,771", "רשומות אוטובוס", TEAL_MID),
    ("508,772", "רשומות רכבת",     TEAL_DARK),
    ("67,423",  "מקטעי תנועה",     ACCENT_GREEN),
    ("51,874",  "עדכוני נסיעה",    TEAL_MID),
    ("17,512",  "התראות שירות",    ACCENT_AMBER),
    ("70+",     "ויזואליזציות",    TEAL_DARK),
]
metrics_active = [
    ("23",  "קווי דן פעילים",      TEAL_MID),
    ("30",  "קווי אגד פעילים",     TEAL_DARK),
    ("15",  "קווי סופרבוס",        ACCENT_GREEN),
]

# First row: total docs (3 columns × 2 rows)
y = Inches(1.45)
for i, (val, label, color) in enumerate(metrics):
    col = i % 3
    row = i // 3
    x = Inches(0.5 + col * 4.25)
    yy = y + row * Inches(1.4)
    box = add_textbox(s, x, yy, Inches(4.0), Inches(1.25),
                      fill=WHITE, line=color)
    # Big number (LTR for digits)
    add_text(s, x, yy + Inches(0.05), Inches(4.0), Inches(0.7),
             val, size=42, bold=True, color=color,
             align=PP_ALIGN.CENTER, rtl=False)
    # Label
    add_text(s, x, yy + Inches(0.75), Inches(4.0), Inches(0.4),
             label, size=14, color=DARK_TXT, align=PP_ALIGN.CENTER)

# Second section title
add_text(s, Inches(0.5), Inches(4.4), Inches(12.3), Inches(0.4),
         "🟢  קווים פעילים עכשיו (velocity > 0):",
         size=16, bold=True, color=TEAL_DARK, align=PP_ALIGN.RIGHT)

# Second row: active lines
for i, (val, label, color) in enumerate(metrics_active):
    x = Inches(0.5 + i * 4.25)
    yy = Inches(4.95)
    box = add_textbox(s, x, yy, Inches(4.0), Inches(1.4),
                      fill=color, line=color)
    add_text(s, x, yy + Inches(0.15), Inches(4.0), Inches(0.7),
             val, size=48, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, rtl=False)
    add_text(s, x, yy + Inches(0.85), Inches(4.0), Inches(0.4),
             label, size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# Footer
add_text(s, Inches(0.5), Inches(6.55), Inches(12.3), Inches(0.4),
         "📡 הנתונים מתעדכנים אוטומטית כל 5 דקות באמצעות cron pipeline",
         size=12, color=GRAY, align=PP_ALIGN.CENTER)


# ─────────── SLIDE 11: Tech Stack ───────────
s = prs.slides.add_slide(blank_layout); slide_bg(s)
add_header(s, "Tech Stack",
           "כל הטכנולוגיות שבשימוש", 11, TOTAL)

categories = [
    ("Backend", TEAL_DARK,
     ["Python 3.12", "Flask (HTTP server)", "asyncio + threading",
      "boto3 (S3/MinIO)", "elasticsearch-py", "kafka-python", "requests"]),
    ("Big Data", TEAL_MID,
     ["Apache Kafka 7.5", "Apache Airflow 2.7", "Elasticsearch 8.8",
      "MinIO (S3-compatible)", "Apache Zookeeper"]),
    ("Database", ACCENT_GREEN,
     ["PostgreSQL 13", "psycopg2", "SQL (GTFS schema)",
      "Schema migrations"]),
    ("Frontend", ACCENT_AMBER,
     ["HTML5 + JavaScript ES6+", "Leaflet (maps)", "Kibana 8.8",
      "Hebrew RTL UI", "Responsive design"]),
    ("DevOps", RED,
     ["Docker + Docker Compose", "Linux VM (Hyper-V)",
      "Bash automation scripts", "Cron pipelines"]),
    ("AI", TEAL_DARK,
     ["Google Gemini API", "Stride GTFS API",
      "HERE Traffic v7 API", "Open Bus Stride"]),
]

y0 = Inches(1.5)
w = Inches(4.15)
h = Inches(2.6)
gap = Inches(0.15)
for i, (cat, color, items) in enumerate(categories):
    col = i % 3
    row = i // 3
    x = Inches(0.4 + col * (4.15 + 0.15))
    y = y0 + row * (h + Inches(0.2))
    box = add_textbox(s, x, y, w, h, fill=WHITE, line=color)
    # Header bar
    head = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
        x, y, w, Inches(0.55))
    head.fill.solid(); head.fill.fore_color.rgb = color
    head.line.fill.background()
    add_text(s, x + Inches(0.15), y + Inches(0.05), w - Inches(0.3), Inches(0.45),
             cat, size=18, bold=True, color=WHITE,
             align=PP_ALIGN.LEFT, rtl=False)
    # Items
    for j, item in enumerate(items):
        add_text(s, x + Inches(0.2), y + Inches(0.65 + j * 0.27),
                 w - Inches(0.4), Inches(0.25),
                 "• " + item, size=11, color=DARK_TXT,
                 align=PP_ALIGN.LEFT, rtl=False)


# ─────────── SLIDE 12: Code Quality / Stats ───────────
s = prs.slides.add_slide(blank_layout); slide_bg(s)
add_header(s, "פרוייקט הקוד  💻",
           "מבנה הקוד, מודולריות ובדיקות", 12, TOTAL)

# Left: directory tree
add_text(s, Inches(0.5), Inches(1.5), Inches(6.0), Inches(0.5),
         "מבנה התיקיות:", size=18, bold=True, color=TEAL_DARK)

tree = """finalproject/
├── bot.py                     # Flask server (2,500 lines)
├── moovit.html                # Frontend UI
├── docker-compose.yml         # 10 services
├── producers/                 # Data producers
│   ├── bus_positions.py
│   ├── train_positions.py
│   └── traffic.py
├── airflow/dags/              # 6 DAGs
├── etl/transformers.py        # Stream ETL
├── storage/
│   ├── minio_uploader.py
│   └── es_indexer.py
├── scripts/
│   ├── create_kibana_dashboards.py
│   ├── watchdog.sh
│   ├── backup_kibana.sh
│   └── demo_prep.sh
└── Gtfs_query.py              # PostgreSQL queries"""

tb = s.shapes.add_textbox(Inches(0.5), Inches(2.0), Inches(6.0), Inches(5.0))
tf = tb.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.LEFT
hebrew_run(p, tree, size=11, color=DARK_TXT, font='Consolas')
# Border
border = add_textbox(s, Inches(0.4), Inches(1.95), Inches(6.2), Inches(5.1),
                     line=TEAL_LIGHT)

# Right: key practices
add_text(s, Inches(7.0), Inches(1.5), Inches(5.9), Inches(0.5),
         "עקרונות תכנון:", size=18, bold=True, color=TEAL_DARK)

practices = [
    ("🔁", "Resilience",
     "ניסיונות חוזרים, fallback בכל שלב — אם API נופל, מערכת ממשיכה"),
    ("🧩", "Modularity",
     "כל producer/transformer בקובץ נפרד, קל להוסיף מקור חדש"),
    ("📦", "Containerization",
     "כל שירות ב-Docker container מבודד עם volumes נפרדים"),
    ("🛡️", "Store-and-Forward",
     "buffer מקומי כשה-MinIO/ES נופלים — flush אוטומטי כשהם חוזרים"),
    ("⏰", "Scheduling",
     "Airflow + cron-pipeline כגיבוי כפול לוודא שהזרם תמיד פעיל"),
    ("🔐", "Configuration",
     ".env לסודות, לא בקוד; פרמטרים נטענים בזמן ריצה"),
    ("📊", "Observability",
     "logs מובנים, /health endpoint, watchdog אוטומטי"),
]
for i, (icon, title, desc) in enumerate(practices):
    y = Inches(2.0 + i * 0.85)
    box = add_textbox(s, Inches(7.0), y, Inches(5.9), Inches(0.78),
                      fill=WHITE, line=TEAL_LIGHT)
    add_text(s, Inches(7.1), y + Inches(0.08), Inches(0.7), Inches(0.7),
             icon, size=24, color=TEAL_DARK,
             align=PP_ALIGN.CENTER, rtl=False)
    add_text(s, Inches(7.8), y + Inches(0.05), Inches(5.0), Inches(0.35),
             title, size=14, bold=True, color=TEAL_DARK,
             align=PP_ALIGN.LEFT, rtl=False)
    add_text(s, Inches(7.8), y + Inches(0.40), Inches(5.0), Inches(0.40),
             desc, size=11, color=DARK_TXT)


# ─────────── SLIDE 13: Challenges Solved ───────────
s = prs.slides.add_slide(blank_layout); slide_bg(s)
add_header(s, "אתגרים שפתרנו  🛠️",
           "תקלות אמיתיות שנתקלנו בהן ופתרנו", 13, TOTAL)

challenges = [
    ("🚫", "MOT GTFS-RT WAF block",
     "פתרון: שימוש ב-Open Bus Stride של עמותת \"הסדנא\" כפרוקסי",
     RED),
    ("🔌", "Container networking",
     "פתרון: הרצת Producers על ה-Host כעוקף עד שתוקנו NAT rules",
     ACCENT_AMBER),
    ("⚠️", "siri_route__ vs siri_routes__",
     "באג נסתר: רבים/יחיד בפרמטרים של API גרם לחזרה של נתונים שגויים",
     RED),
    ("📉", "Bogus 2038 train timestamps",
     "פתרון: סינון range filter במקום sort על recorded_at",
     ACCENT_AMBER),
    ("🐌", "Elasticsearch bulk timeouts",
     "פתרון: chunking ל-50 docs בכל batch + retry עם exponential backoff",
     TEAL_MID),
    ("📋", "Duplicate stop names",
     "פתרון: מעבר מ-siri_ride_stops ל-gtfs_ride_stops עם dedupe לפי sequence",
     ACCENT_GREEN),
]

y = Inches(1.45)
h = Inches(0.85)
for icon, title, sol, color in challenges:
    box = add_textbox(s, Inches(0.4), y, Inches(12.5), h,
                      fill=WHITE, line=color)
    # Color stripe
    stripe = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
        Inches(12.6), y, Inches(0.3), h)
    stripe.fill.solid(); stripe.fill.fore_color.rgb = color
    stripe.line.fill.background()
    # Icon
    add_text(s, Inches(0.6), y + Inches(0.15), Inches(0.8), Inches(0.6),
             icon, size=28, color=color, align=PP_ALIGN.CENTER, rtl=False)
    # Title
    add_text(s, Inches(1.5), y + Inches(0.1), Inches(11.0), Inches(0.4),
             title, size=15, bold=True, color=DARK_TXT,
             align=PP_ALIGN.LEFT, rtl=False)
    # Solution
    add_text(s, Inches(1.5), y + Inches(0.45), Inches(11.0), Inches(0.4),
             sol, size=12, color=GRAY)
    y += h + Inches(0.05)


# ─────────── SLIDE 14: Future Work ───────────
s = prs.slides.add_slide(blank_layout); slide_bg(s)
add_header(s, "המשך פיתוח  🚀",
           "מה אפשר להוסיף בעתיד", 14, TOTAL)

# 6 future ideas in 2 rows × 3 cols
ideas = [
    ("🤖", "ML Predictions",
     "מודל ML לחיזוי איחורים על בסיס היסטוריה + תנאי תנועה"),
    ("📱", "Mobile App",
     "React Native PWA או אפליקציה native ל-iOS/Android"),
    ("🗺️", "Heat Maps",
     "מפת חום של עומסי תנועה לפי שעה ביום + יום בשבוע"),
    ("🔔", "Real-time Alerts",
     "התראות Push למשתמשים על שיבושים בקווים שלהם"),
    ("📈", "Historical Analytics",
     "מחסן נתונים ב-Redshift/BigQuery לניתוחים ארוכי טווח"),
    ("🌍", "All-Israel Coverage",
     "הרחבה מגוש דן לכל הארץ + אינטגרציה עם תחבורה בערים אחרות"),
]
w = Inches(4.15)
h = Inches(2.45)
gap = Inches(0.15)
y0 = Inches(1.7)
for i, (icon, title, desc) in enumerate(ideas):
    col = i % 3
    row = i // 3
    x = Inches(0.4 + col * (4.15 + 0.15))
    y = y0 + row * (h + Inches(0.25))
    box = add_textbox(s, x, y, w, h, fill=WHITE, line=TEAL_LIGHT)
    # Icon big at top
    add_text(s, x, y + Inches(0.2), w, Inches(0.9),
             icon, size=54, color=TEAL_MID,
             align=PP_ALIGN.CENTER, rtl=False)
    # Title
    add_text(s, x + Inches(0.15), y + Inches(1.15), w - Inches(0.3), Inches(0.4),
             title, size=18, bold=True, color=TEAL_DARK,
             align=PP_ALIGN.CENTER, rtl=False)
    # Desc
    add_text(s, x + Inches(0.2), y + Inches(1.6), w - Inches(0.4), Inches(0.8),
             desc, size=12, color=DARK_TXT, align=PP_ALIGN.CENTER)


# ─────────── SLIDE 15: Thank You ───────────
s = prs.slides.add_slide(blank_layout); slide_bg(s, TEAL_LIGHT)

# Decorative circles
for x, y, r, c in [(2.5, 1.5, 1.5, SKY_BLUE), (11.5, 6.0, 1.2, TEAL_MID),
                   (10.8, 1.0, 0.8, TEAL_MID), (1.5, 6.5, 1.0, SKY_BLUE)]:
    c1 = s.shapes.add_shape(MSO_SHAPE.OVAL,
        Inches(x - r), Inches(y - r), Inches(2*r), Inches(2*r))
    c1.fill.solid(); c1.fill.fore_color.rgb = c
    c1.line.fill.background()

add_text(s, Inches(0.5), Inches(2.5), Inches(12.3), Inches(1.2),
         "תודה רבה!  🙏",
         size=84, bold=True, color=TEAL_DARK, align=PP_ALIGN.CENTER)

add_text(s, Inches(0.5), Inches(3.9), Inches(12.3), Inches(0.6),
         "שאלות?",
         size=36, color=DARK_TXT, align=PP_ALIGN.CENTER)

add_text(s, Inches(0.5), Inches(5.0), Inches(12.3), Inches(0.5),
         "🚌  ניטור תחבורה ציבורית בזמן אמת — גוש דן ותל אביב",
         size=18, color=TEAL_DARK, align=PP_ALIGN.CENTER)

add_text(s, Inches(0.5), Inches(5.6), Inches(12.3), Inches(0.4),
         "Israel Public Transit Real-Time Monitoring Platform",
         size=13, color=GRAY, align=PP_ALIGN.CENTER, rtl=False, font='Arial')


# ─────────── Save ───────────
import os
out_path = r"C:\Users\User\Downloads\Israel_Transit_Project_Presentation.pptx"
prs.save(out_path)
print(f"Saved: {out_path}")
print(f"Slides: {len(prs.slides)}")
print(f"Size: {os.path.getsize(out_path):,} bytes")
